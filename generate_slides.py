"""
Automated PowerPoint Presentation Generator
Creates a basic slide deck structure for the waste management project.

Usage:
    python generate_slides.py

Note: You'll need to add the chart images manually or run export_charts.py first.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

print("=" * 80)
print("AUTOMATED POWERPOINT GENERATOR - WASTE MANAGEMENT PROJECT")
print("=" * 80)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define paths
assets_dir = Path(__file__).parent / 'assets'

# ==================== SLIDE 1: TITLE ====================
print("\n1. Creating Title Slide...")

slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "ENVIRONMENTAL DASHBOARD"
subtitle.text = "Comparative Waste Management Analysis\nEurope & Africa (1990-2021)\n\nBellatreche Mohamed Amine\nCherif Ghizlane Imane\nNovembre 2025"

# Format title
title_frame = title.text_frame
title_frame.paragraphs[0].font.size = Pt(44)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(31, 119, 180)

# Format subtitle
subtitle_frame = subtitle.text_frame
subtitle_frame.paragraphs[0].font.size = Pt(20)

print("   ✓ Title slide created")

# ==================== SLIDE 2: CONTEXT & OBJECTIVES ====================
print("2. Creating Context & Objectives Slide...")

slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide2.shapes.title
title.text = "Context & Project Objectives"

# Add text box for content
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(5.5)

txBox = slide2.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True

# Add context
p1 = tf.paragraphs[0]
p1.text = "🌍 GLOBAL CONTEXT"
p1.font.size = Pt(20)
p1.font.bold = True
p1.space_after = Pt(10)

p2 = tf.add_paragraph()
p2.text = "• 49 countries analyzed (27 Europe + 22 Africa)"
p2.font.size = Pt(16)
p2.level = 1

p3 = tf.add_paragraph()
p3.text = "• North-South divide in waste management infrastructure"
p3.font.size = Pt(16)
p3.level = 1

p4 = tf.add_paragraph()
p4.text = "• Europe: 25-35% recycling rates, advanced systems"
p4.font.size = Pt(16)
p4.level = 1

p5 = tf.add_paragraph()
p5.text = "• Africa: Rapid waste growth (3-5%/year), limited infrastructure"
p5.font.size = Pt(16)
p5.level = 1
p5.space_after = Pt(20)

# Add objectives
p6 = tf.add_paragraph()
p6.text = "🎯 OBJECTIVES"
p6.font.size = Pt(20)
p6.font.bold = True
p6.space_after = Pt(10)

p7 = tf.add_paragraph()
p7.text = "• Comparative analysis: Europe vs. Africa waste practices"
p7.font.size = Pt(16)
p7.level = 1

p8 = tf.add_paragraph()
p8.text = "• Machine learning forecasts (ARIMA models)"
p8.font.size = Pt(16)
p8.level = 1

p9 = tf.add_paragraph()
p9.text = "• Environmental risk assessment (rule-based scoring)"
p9.font.size = Pt(16)
p9.level = 1

p10 = tf.add_paragraph()
p10.text = "• Evidence-based policy recommendations"
p10.font.size = Pt(16)
p10.level = 1

print("   ✓ Context & Objectives slide created")

# ==================== SLIDE 3: DATA SOURCES ====================
print("3. Creating Data Sources Slide...")

slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide3.shapes.title
title.text = "Data Sources & Methodology"

# Add content
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(5.5)

txBox = slide3.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "📊 DATA SOURCES"
p1.font.size = Pt(20)
p1.font.bold = True
p1.space_after = Pt(10)

p2 = tf.add_paragraph()
p2.text = "• Waste generation: UN Environment (132 countries, 2000-2021)"
p2.font.size = Pt(16)
p2.level = 1

p3 = tf.add_paragraph()
p3.text = "• Recycling rates: OECD Municipal Waste (38 countries, 1990-2015)"
p3.font.size = Pt(16)
p3.level = 1

p4 = tf.add_paragraph()
p4.text = "• Linear interpolation: Biennial data → Annual time series"
p4.font.size = Pt(16)
p4.level = 1
p4.space_after = Pt(20)

p5 = tf.add_paragraph()
p5.text = "🎯 KEY INDICATORS"
p5.font.size = Pt(20)
p5.font.bold = True
p5.space_after = Pt(10)

indicators = [
    "Recycling rate (%, Europe only)",
    "Waste per capita (kg/year)",
    "Waste by sector (households, construction, manufacturing, services)",
    "Risk score (0-100, rule-based)",
    "5-year ARIMA forecasts"
]

for ind in indicators:
    p = tf.add_paragraph()
    p.text = f"✓ {ind}"
    p.font.size = Pt(16)
    p.level = 1

print("   ✓ Data Sources slide created")

# ==================== SLIDE 4: RESULTS - TRENDS ====================
print("4. Creating Results Slide (Trends)...")

slide4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "Waste Management: Europe vs. Africa (1990-2021)"
title_para.font.size = Pt(28)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

# Check if chart exists
chart_path = assets_dir / 'chart1_total_waste_trend.png'
if chart_path.exists():
    slide4.shapes.add_picture(
        str(chart_path),
        Inches(0.5),
        Inches(1.2),
        width=Inches(9)
    )
    print("   ✓ Results slide created with chart")
else:
    # Add placeholder text
    placeholder = slide4.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
    tf = placeholder.text_frame
    p = tf.paragraphs[0]
    p.text = "[Insérer graphique: chart1_total_waste_trend.png]\n\nExécutez export_charts.py pour générer les graphiques"
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    print("   ⚠ Results slide created (placeholder - run export_charts.py)")

# Add key metrics
metrics_box = slide4.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(1))
tf = metrics_box.text_frame
p = tf.paragraphs[0]
p.text = "🇪🇺 Europe: 500 kg/cap/yr avg, 30% recycling  |  🌍 Africa: 250 kg/cap/yr avg, 3-5% growth"
p.font.size = Pt(16)
p.alignment = PP_ALIGN.CENTER

# ==================== SLIDE 5: COMPOSITION ====================
print("5. Creating Composition Slide...")

slide5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "Waste Composition by Economic Sector"
title_para.font.size = Pt(28)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

# Check if charts exist
chart5_path = assets_dir / 'chart5_composition_pie.png'
chart7_path = assets_dir / 'chart7_composition_bar.png'

if chart5_path.exists():
    slide5.shapes.add_picture(
        str(chart5_path),
        Inches(0.5),
        Inches(1.2),
        width=Inches(4.5)
    )

if chart7_path.exists():
    slide5.shapes.add_picture(
        str(chart7_path),
        Inches(5),
        Inches(1.2),
        width=Inches(4.5)
    )

if not chart5_path.exists() or not chart7_path.exists():
    placeholder = slide5.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(1.5))
    tf = placeholder.text_frame
    p = tf.paragraphs[0]
    p.text = "[Insérer graphiques de composition]\n\nExécutez export_charts.py"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    print("   ⚠ Composition slide created (placeholder)")
else:
    print("   ✓ Composition slide created with charts")

# Add insights
insights_box = slide5.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(1))
tf = insights_box.text_frame
p = tf.paragraphs[0]
p.text = "✓ Households = largest sector | ✓ Construction waste significant | ✓ Stacked area shows temporal evolution"
p.font.size = Pt(14)
p.alignment = PP_ALIGN.CENTER

# ==================== SLIDE 6: RECOMMENDATIONS ====================
print("6. Creating Recommendations Slide...")

slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide6.shapes.title
title.text = "Strategic Recommendations"

# Add content
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(5.5)

txBox = slide6.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True

# Recommendations in 3 columns (simulated)
p1 = tf.paragraphs[0]
p1.text = "🇪🇺 EUROPE (OPTIMIZATION)"
p1.font.size = Pt(18)
p1.font.bold = True
p1.space_after = Pt(5)

recs1 = [
    "Circular economy transition",
    "Meet 50% recycling targets",
    "Extended Producer Responsibility",
    "Harmonize data collection"
]

for rec in recs1:
    p = tf.add_paragraph()
    p.text = f"• {rec}"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(3)

p_sep1 = tf.add_paragraph()
p_sep1.space_after = Pt(10)

p2 = tf.add_paragraph()
p2.text = "🌍 AFRICA (INFRASTRUCTURE)"
p2.font.size = Pt(18)
p2.font.bold = True
p2.space_after = Pt(5)

recs2 = [
    "Basic collection systems",
    "Formalize informal sector",
    "Regional cooperation",
    "Technology transfer from Europe"
]

for rec in recs2:
    p = tf.add_paragraph()
    p.text = f"• {rec}"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(3)

p_sep2 = tf.add_paragraph()
p_sep2.space_after = Pt(10)

p3 = tf.add_paragraph()
p3.text = "🤖 ML & ANALYTICS"
p3.font.size = Pt(18)
p3.font.bold = True
p3.space_after = Pt(5)

recs3 = [
    "ARIMA forecasting: 5-year trends",
    "Risk assessment: Priority countries",
    "Scenario modeling: Policy impact",
    "Dashboard for evidence-based decisions"
]

for rec in recs3:
    p = tf.add_paragraph()
    p.text = f"• {rec}"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(3)

print("   ✓ Recommendations slide created")

# ==================== SLIDE 7: THANK YOU ====================
print("7. Creating Thank You Slide...")

slide7 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

# Title
title_box = slide7.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
title_frame = title_box.text_frame
title_para = title_frame.paragraphs[0]
title_para.text = "❓ QUESTIONS ?"
title_para.font.size = Pt(44)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

# Thank you
thanks_box = slide7.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(0.8))
thanks_frame = thanks_box.text_frame
thanks_para = thanks_frame.paragraphs[0]
thanks_para.text = "Merci de votre attention"
thanks_para.font.size = Pt(28)
thanks_para.alignment = PP_ALIGN.CENTER

# Contact
contact_box = slide7.shapes.add_textbox(Inches(2), Inches(4.8), Inches(6), Inches(1))
contact_frame = contact_box.text_frame
contact_para = contact_frame.paragraphs[0]
contact_para.text = "📧 [votre.email@institution.dz]\n📊 Dashboard disponible sur demande"
contact_para.font.size = Pt(16)
contact_para.alignment = PP_ALIGN.CENTER

# Sources
sources_box = slide7.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
sources_frame = sources_box.text_frame
sources_para = sources_frame.paragraphs[0]
sources_para.text = "─────────────────────────────\nSources: UN Environment (132 countries) | OECD Municipal Waste (38 countries) | World Bank Population"
sources_para.font.size = Pt(12)
sources_para.alignment = PP_ALIGN.CENTER
sources_para.font.color.rgb = RGBColor(128, 128, 128)

print("   ✓ Thank You slide created")

# ==================== SAVE PRESENTATION ====================
output_path = Path(__file__).parent / 'presentation_waste_management_algeria.pptx'
prs.save(output_path)

print("\n" + "=" * 80)
print("PRESENTATION CREATED SUCCESSFULLY!")
print("=" * 80)
print(f"\n✓ File saved: {output_path.absolute()}")
print("\n📊 Slides created:")
print("  1. Title slide")
print("  2. Context & Objectives")
print("  3. Data Sources & Methodology")
print("  4. Results - Trends (with chart placeholder)")
print("  5. Composition (with chart placeholders)")
print("  6. Recommendations")
print("  7. Thank You / Questions")

print("\n⚠️ IMPORTANT NEXT STEPS:")
print("  1. Run: python export_charts.py")
print("     → This will generate high-quality chart images")
print("\n  2. Open: presentation_waste_management_algeria.pptx")
print("     → Replace placeholders with actual charts from assets/")
print("\n  3. Customize:")
print("     → Add your name and institution")
print("     → Adjust colors/theme if desired")
print("     → Fill in the XX% placeholders with actual numbers")
print("     → Add your email on the last slide")

print("\n💡 TIPS:")
print("  - Charts are in assets/ folder after running export_charts.py")
print("  - You can change the theme: Design tab → Themes")
print("  - Add animations: Animations tab (optional)")
print("  - Test your presentation in Slide Show mode")

print("\n✅ Your presentation structure is ready!")
print("=" * 80)
